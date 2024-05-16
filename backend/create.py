import pptx.presentation
import base64
from pptx.util import Inches, Pt

Title_Font_Size = Pt(30)
Slide_Font_Size = Pt(16)


def generate_slide_title(topic):
    prompt = f"Generate 5 slides title for the given topic '{topic}' "
    response = client.completions.create(
        model="GPT-3.5 Turbo", prompt=prompt, max_tokens=200
    )
    return response.choices[0].text.split("\n")


def generate_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'"
    response = client.completions.create(
        model="GPT-3.5 Turbo", prompt=prompt, max_tokens=500
    )
    return response.choices[0].text


@app.post("/ppt")
async def create_presentation(topic: str):
    try:
        slide_titles = generate_slide_title(topic)
        slide_contents = [generate_slide_content(title) for title in slide_titles]
        prs = pptx.Presentation()
        slide_layout = prs.slide_layout[1]
        title_slide = prs.slides.add_slides(prs.slide_layout[0])
        title_slide.shapes.title.text = topic

        for slide_title, slide_content in zip(slide_titles, slide_contents):
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_title
            slide.shapes.placeholders[1].text = slide_content

            slide.shapes.title.text_frame.paragraphs[0].font_size = Title_Font_Size
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        paragraph.font_size = Slide_Font_Size

        output_path = f"generated_ppt/{topic}_presentation.pptx"
        prs.save(output_path)
        return {"file_path": output_path}
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error generating presentation: {str(e)}"
        )
