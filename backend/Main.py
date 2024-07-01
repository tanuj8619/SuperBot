import base64  # For encoding and decoding base64 strings (used for audio data).
import csv
from io import BytesIO, StringIO  # For handling in-memory binary and text streams.
import io
from typing import Optional
from fastapi import (
    FastAPI,
    Form,
    HTTPException,
    Query,
    UploadFile,
    File,
)  # Main FastAPI class to create the app.
from fastapi.responses import JSONResponse, FileResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel  # For defining data models.
from openai import AzureOpenAI
from gtts import gTTS
import google.generativeai as genai
from config import AZURE_OPENAI_API_KEY, GOOGLE_API_KEY
import PyPDF2
import docx
import pptx.presentation
from pptx.util import Inches, Pt
from pptx import Presentation
from docx import Document
import re  # For regular expression operations.
import os  # For interacting with the operating system.

app = FastAPI()  # instance of the FastAPI class to set up the web application.

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Update this with your frontend URL
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE"],
    allow_headers=["*"],
)

genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel("gemini-pro")
chats = model.start_chat(
    history=[]
)  # Initializes a chat session with an empty history.


class UserInput(BaseModel):
    user_input: str


class Gemini(BaseModel):
    gemini: bool


class ConversationContext:
    def __init__(self):
        self.history = []
        # self.latest_chat = {}


context = ConversationContext()

client = AzureOpenAI(
    azure_endpoint="https://pss-openai-service.openai.azure.com/",
    api_key=AZURE_OPENAI_API_KEY,
    api_version="2024-02-01",
)


def extract_text_from_pdf(file_content):
    pdf_reader = PyPDF2.PdfFileReader(file_content)
    text = ""
    for page_num in range(len(pdf_reader.pages)):
        text += pdf_reader.pages[page_num].extract_text()
        text += "\n"  # Add a newline after each page
    return text


def extract_text_from_docx(docx_file):
    doc = docx.Document(docx_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"  # Add a newline after each paragraph
    return text


def extract_text_from_txt(txt_file):
    return txt_file.read().decode("utf-8")


# Title_Font_Size = Pt(30)
# Slide_Font_Size = Pt(16)


# def generate_slide_title(topic):
#     prompt = f"Generate 5 slides title for the given topic '{topic}' "
#     response = model.generate_content(prompt)
#     return response.text.split("\n")


# def generate_slide_content(slide_title):
#     prompt = f"Generate content for the slide: '{slide_title}'"
#     response = model.generate_content(prompt)
#     return response.text


def get_gemini_response(prompt):
    model = genai.GenerativeModel("gemini-pro")
    response = model.generate_content(prompt)
    return response


def refine_subtopics(sub_topics):  # removing the first 3 characters and quotes
    sub_titles = []
    for sub_topic in sub_topics:
        sub_titles.append(sub_topic[3:].replace('"', ""))
    return sub_titles


def content_generation(sub_titles):
    content = []
    for i in sub_titles:
        prompt = f"Generate a content of {i} for presentation slide on the 2 bullet point only each of point 20 tokens"
        model = genai.GenerativeModel("gemini-pro")
        response = model.generate_content(prompt)
        content.append(response.text)
    return content


def refine_final_content(content):
    final_content = []
    for i in content:
        cleaned_text = clean_text(i)
        sentences = split_sentences(cleaned_text)
        final_content.append(sentences)
    print("final content ready....")
    return final_content


def clean_text(text):
    # Remove extra whitespaces and newlines
    cleaned_text = re.sub("\s+", " ", text).strip()

    # Remove markdown-style bullet points, asterisks, and numeric bullet points
    cleaned_text = re.sub(r"[*-]\s*|\d+\.\s*", "", cleaned_text)

    # Remove extra spaces before and after colons
    cleaned_text = re.sub(r"\s*:\s*", ": ", cleaned_text)

    # Remove extra spaces before and after hyphens
    cleaned_text = re.sub(r"\s*-\s*", " - ", cleaned_text)

    return cleaned_text


def split_sentences(text):
    # Split the text into sentences using regular expression
    sentences = re.split(r"(?<=\.)\s+", text)

    # Capitalize the first letter of each sentence
    sentences = [sentence.capitalize() for sentence in sentences]

    return sentences


def replace_and_capitalize(text):
    # Define a function to replace and capitalize the text between colons
    def replace_and_capitalize_colon(match):
        return match.group(1) + match.group(2).capitalize() + match.group(3)

    # Use regular expression to find and replace text between colons
    result = re.sub(r"(:\s*)(.*?)(\s*:[^:]|$)", replace_and_capitalize_colon, text)

    return result


def slide_maker(powerpoint, topic, sub_titles, final_content):
    title_slide_layout = powerpoint.slide_layouts[0]
    title_slide = powerpoint.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    content = title_slide.placeholders[1]
    content.text = "Created By AI Model"
    for i in range(len(sub_titles)):
        bulletLayout = powerpoint.slide_layouts[1]
        secondSlide = powerpoint.slides.add_slide(bulletLayout)
        # accessing the attributes of shapes
        myShapes = secondSlide.shapes
        titleShape = myShapes.title
        bodyShape = myShapes.placeholders[1]
        titleShape.text = sub_titles[i]
        titleShape.text_frame.paragraphs[0].font.size = Pt(24)
        titleShape.text_frame.paragraphs[0].font.bold = True
        tFrame = bodyShape.text_frame
        print("Topic Generated")
        for point in final_content[i]:
            point = re.sub(r":[^:]+:", ":", point)
            point = replace_and_capitalize(point)
            p = tFrame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.space_after = Pt(10)
    return powerpoint


def generate_csv_content(prompt):
    # Generate structured content based on the prompt
    model = genai.GenerativeModel("gemini-pro")
    response = model.generate_content(prompt)
    print(response)

    # Replace asterisks with empty string in the generated response
    text_content = response.text.replace("*", "")

    # Assuming the generated response contains rows of data
    csv_data = text_content.strip().split("\n")

    if len(csv_data) == 0:
        raise ValueError("No data generated")

    headers = csv_data[0].split(",")  # Assume the first line contains headers
    rows = [row.split(",") for row in csv_data[1:]]  # Remaining lines are data rows

    output = StringIO()
    writer = csv.writer(output)
    writer.writerow(headers)
    writer.writerows(rows)

    return output.getvalue()


# Function to remove asterisks around text
def remove_asterisks(text):
    # Remove asterisks at the beginning and end of each line, considering Markdown
    lines = text.splitlines()
    cleaned_lines = []
    for line in lines:
        line = re.sub(r"^\s*[#\*]+\s*", "", line)  # Remove asterisks at the beginning
        line = re.sub(r"\s*[#\*]+\s*$", "", line)  # Remove asterisks at the end
        line = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
        cleaned_lines.append(line.strip())
    return "\n".join(cleaned_lines)


@app.post("/doc")
async def create_word_document(topic: str = Form(None)):
    prompt = f"{topic} "

    response = get_gemini_response(prompt)
    print("content Generated")
    print(response.text)
    doc = Document()
    doc.add_heading(remove_asterisks(topic), 0)
    doc.add_paragraph(remove_asterisks(response.text))

    # Save the document to a BytesIO object
    doc_bytes = io.BytesIO()
    doc.save(doc_bytes)
    doc_bytes.seek(0)  # Reset the BytesIO position to the beginning

    # Return the document file as a streaming response
    file_name = f"{topic}_document.docx"
    return StreamingResponse(
        iter([doc_bytes.getvalue()]),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={file_name}"},
    )

    # output_path = f"generated_doc/{topic}_document.docx"
    # doc.save(output_path)
    # return {"file_path": output_path}


@app.post("/csv")
async def generate_csv(topic: str = Form(None)):
    try:
        if not topic:
            raise HTTPException(status_code=400, detail="User input is required")

        prompt = f"Generate data relevant to the: {topic}"
        csv_content = generate_csv_content(prompt)
        csv_bytes = BytesIO(csv_content.encode("utf-8"))
        csv_bytes.seek(0)

        file_name = "generated_file.csv"
        return StreamingResponse(
            iter([csv_bytes.getvalue()]),
            media_type="text/csv",
            headers={"Content-Disposition": f"attachment; filename={file_name}"},
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating CSV: {str(e)}")


@app.post("/ppt")
async def create_presentation(topic: str = Form(None)):
    try:
        prompt = f"Generate a 5 sub-titles only  on the topic of {topic}"
        response = get_gemini_response(prompt)
        print("Topic Generated")
        sub_topics = response.text.split("\n")
        sub_titles = refine_subtopics(sub_topics)
        print("Sub Titles")
        content = content_generation(sub_titles)
        print("content Generated")
        final_content = refine_final_content(content)
        print("final content ready")
        powerpoint = Presentation()
        powerpoint = slide_maker(powerpoint, topic, sub_titles, final_content)
        print("presenatation ready:")

        # Save the presentation to a BytesIO(buffer) object
        presentation_bytes = io.BytesIO()
        powerpoint.save(presentation_bytes)
        print("presenatation ready:")
        presentation_bytes.seek(0)  # Reset the BytesIO position to the beginning

        # Return the presentation file as a streaming response
        file_name = f"{topic}_presentation.pptx"
        # StreamingResponse used in web frameworks (like FastAPI) to stream large files to the client without loading the entire file into memory at once.
        return StreamingResponse(
            iter([presentation_bytes.getvalue()]),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",  # Sets the MIME type of the response to the appropriate type for a PowerPoint file .pptx
            headers={
                "Content-Disposition": f"attachment; filename={file_name}"
            },  # This header tells the browser to download the file and use the provided filename.
        )
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error generating presentation: {str(e)}"
        )


@app.post("/api/chat")
async def chat(
    user_input: Optional[str] = Form(None),
    file: UploadFile = File(None),
    gemini: Optional[bool] = Form(None),
):
    print("User Input:", user_input)
    print("isgenai:", gemini)
    extracted_text = " "

    if not gemini:
        try:
            if user_input is None and file is None:
                raise HTTPException(
                    status_code=400,
                    detail="Either user input or file upload is required.",
                )

            if file is not None:
                # Extract text from uploaded file
                file_extension = file.filename.split(".")[-1]
                if file_extension == "pdf":
                    input_text = extract_text_from_pdf(file.file)
                elif file_extension == "docx":
                    input_text = extract_text_from_docx(file.file)
                elif file_extension == "txt":
                    input_text = extract_text_from_txt(file.file)
                else:
                    raise HTTPException(
                        status_code=400, detail="Unsupported file format"
                    )
                context.history.append({"role": "user", "content": input_text})
                extracted_text = input_text
                if user_input is not None:
                    context.history.append({"role": "user", "content": user_input})
            elif user_input is not None:
                # User provides input text
                context.history.append({"role": "user", "content": user_input})
            else:
                # No user input provided and no file uploaded, retain previous input text
                if len(context.history) == 0:
                    raise HTTPException(
                        status_code=400, detail="No previous input available."
                    )
                input_text = context.history[-1]["content"]

            # Filter the conversation history based on the context
            filtered_history = [
                message for message in context.history if message["role"] == "user"
            ]

            response = client.chat.completions.create(
                model="gpt-35-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant."},
                    *filtered_history,  # Include filtered conversation history
                ],
            )
            token_used = response.usage.total_tokens
            print(token_used)
            # Append AI response to conversation history
            context.history.append(
                {"role": "assistant", "content": response.choices[0].message.content}
            )

            tts = gTTS(
                text=response.choices[0].message.content, lang="en"
            )  # Specify language as needed
            audio_file = BytesIO()
            tts.write_to_fp(audio_file)
            audio_file.seek(0)

            # Encode audio data as Base64
            audio_data_base64 = base64.b64encode(audio_file.read()).decode("utf-8")

            return {
                "response": response.choices[0].message.content,
                "token_used": token_used,
                "extracted_text": extracted_text,
                "audio": audio_data_base64,
            }

        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

    elif gemini:
        try:
            if user_input is None and file is None:
                raise HTTPException(
                    status_code=400,
                    detail="Either user input or file upload is required.",
                )

            # Initialize extracted_text to None
            extracted_text = None

            if file is not None:
                file_extension = file.filename.split(".")[-1]
                if file_extension == "pdf":
                    extracted_text = extract_text_from_pdf(file.file)
                elif file_extension == "docx":
                    extracted_text = extract_text_from_docx(file.file)
                elif file_extension == "txt":
                    extracted_text = extract_text_from_txt(file.file)
                else:
                    raise HTTPException(
                        status_code=400, detail="Unsupported file format"
                    )

            if user_input is not None:
                # If both file and user_input are provided, send both to the LLM model
                if extracted_text is not None:
                    # Both file extraction and user input are available
                    response = chats.send_message(f"{extracted_text} {user_input}")
                else:
                    # Only user input is available
                    response = chats.send_message(user_input)

                tts = gTTS(text=response.text, lang="en")  # Specify language as needed
                audio_file = BytesIO()
                tts.write_to_fp(audio_file)
                audio_file.seek(0)

                # Encode audio data as Base64
                audio_data_base64 = base64.b64encode(audio_file.read()).decode("utf-8")

                return {
                    "response": response.text,
                    "extracted_text": extracted_text,
                    "audio": audio_data_base64,
                }

            else:
                # No user input provided
                if extracted_text is not None:
                    response = chats.send_message(extracted_text)
                    tts = gTTS(
                        text=response.text, lang="en"
                    )  # Specify language as needed
                    audio_file = BytesIO()
                    tts.write_to_fp(audio_file)
                    audio_file.seek(0)

                    # Encode audio data as Base64
                    audio_data_base64 = base64.b64encode(audio_file.read()).decode(
                        "utf-8"
                    )
                    return {
                        "response": response.text,
                        "extracted_text": extracted_text,
                        "audio": audio_data_base64,
                    }
                else:
                    # No file or user input provided
                    return {"response": "No user message found"}

        except Exception as e:
            raise HTTPException(
                status_code=500, detail="Error processing input: " + str(e)
            )


# for debugging with breakpoints

# if __name__== '__main__':
#     uvicorn.run(app,host='0.0.0.0',port=8000)
